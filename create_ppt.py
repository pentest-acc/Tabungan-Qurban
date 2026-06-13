"""
Jalankan: pip install python-pptx pillow && python create_ppt.py
Output  : Sistem_Tabungan_Qurban.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import datetime

# ── Palet warna ──────────────────────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x0F, 0x17, 0x2A)   # biru gelap
C_PRIMARY    = RGBColor(0x16, 0xA3, 0x4A)   # hijau utama
C_PRIMARY_LT = RGBColor(0x4A, 0xDE, 0x80)   # hijau terang
C_ACCENT     = RGBColor(0x06, 0xB6, 0xD4)   # teal/cyan
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
C_MID_GRAY   = RGBColor(0x94, 0xA3, 0xB8)
C_TEXT_DARK  = RGBColor(0x1E, 0x29, 0x3B)
C_CARD_BG    = RGBColor(0x1E, 0x29, 0x3B)
C_YELLOW     = RGBColor(0xFB, 0xBF, 0x24)
C_PLACEHOLDER= RGBColor(0xE2, 0xE8, 0xF0)   # abu muda untuk placeholder

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # blank layout

# ── Utilitas ──────────────────────────────────────────────────────────────────

def rgb_hex(r, g, b):
    return RGBColor(r, g, b)

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(14), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = font_size
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_img_placeholder(slide, l, t, w, h, label="[ Sisipkan Gambar Di Sini ]",
                        bg=C_PLACEHOLDER, border=C_MID_GRAY, label_color=C_MID_GRAY,
                        sub_label=""):
    """Kotak abu-abu dengan teks label di tengah — tempat diagram."""
    box = add_rect(slide, l, t, w, h, fill=bg, line=border, line_w=Pt(1.5))
    # teks utama
    mid_t = t + h // 2 - Inches(0.35)
    add_text(slide, label, l, mid_t, w, Inches(0.45),
             font_size=Pt(14), bold=True, color=label_color, align=PP_ALIGN.CENTER)
    if sub_label:
        add_text(slide, sub_label, l, mid_t + Inches(0.45), w, Inches(0.35),
                 font_size=Pt(11), color=label_color, align=PP_ALIGN.CENTER)
    return box

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_divider(slide, t, color=C_PRIMARY, w_frac=0.1):
    bar_w = Inches(13.33 * w_frac)
    add_rect(slide, Inches(0.6), t, bar_w, Pt(4), fill=color)

def slide_header(slide, title, subtitle="", dark=True):
    """Header baris atas dengan judul & subtitle."""
    bg_color = C_DARK_BG if dark else C_WHITE
    set_slide_bg(slide, bg_color)
    txt_color = C_WHITE if dark else C_TEXT_DARK

    # strip dekoratif kiri
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill=C_PRIMARY)

    # judul
    add_text(slide, title,
             Inches(0.6), Inches(0.28), Inches(11.5), Inches(0.72),
             font_size=Pt(26), bold=True, color=txt_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.6), Inches(0.95), Inches(11), Inches(0.45),
                 font_size=Pt(13), color=C_MID_GRAY if dark else rgb_hex(0x64,0x74,0x8B),
                 align=PP_ALIGN.LEFT)
    # garis bawah header
    add_rect(slide, Inches(0.6), Inches(1.35), Inches(12.1), Pt(2), fill=C_PRIMARY)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ═════════════════════════════════════════════════════════════════════════════
def slide_cover(prs):
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, C_DARK_BG)

    # blok hijau kiri
    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=C_PRIMARY)

    # titik aksen kanan atas
    add_rect(slide, Inches(10.5), Inches(0), Inches(2.83), Inches(2.2),
             fill=rgb_hex(0x16, 0x29, 0x3E))

    # logo placeholder
    add_img_placeholder(slide,
                        Inches(11.2), Inches(0.2), Inches(1.7), Inches(1.7),
                        label="[ LOGO ]", sub_label="", bg=rgb_hex(0x1E,0x29,0x3B),
                        border=C_PRIMARY_LT, label_color=C_PRIMARY_LT)

    # tag
    add_text(slide, "LAPORAN PROYEK SISTEM INFORMASI",
             Inches(0.55), Inches(1.3), Inches(9), Inches(0.5),
             font_size=Pt(12), color=C_ACCENT, bold=True)

    # judul besar
    add_text(slide, "SISTEM TABUNGAN QURBAN",
             Inches(0.55), Inches(2.0), Inches(10), Inches(1.1),
             font_size=Pt(44), bold=True, color=C_WHITE)
    add_text(slide, "Berbasis Web — React · Node.js · MongoDB",
             Inches(0.55), Inches(3.05), Inches(10), Inches(0.55),
             font_size=Pt(18), color=C_PRIMARY_LT)

    # garis
    add_rect(slide, Inches(0.55), Inches(3.65), Inches(5), Pt(3), fill=C_PRIMARY)

    # deskripsi singkat
    desc = ("Platform digital pengelolaan tabungan qurban masjid — "
            "memungkinkan jamaah menyetor secara bertahap, admin memantau kelompok, "
            "dan sistem pembayaran berjalan otomatis.")
    add_text(slide, desc,
             Inches(0.55), Inches(3.8), Inches(9.5), Inches(1.1),
             font_size=Pt(13), color=C_MID_GRAY)

    # kotak info bawah
    infos = [
        ("Tahun", str(datetime.date.today().year)),
        ("Platform", "Web App (SPA)"),
        ("Stack", "React + Express + MongoDB"),
    ]
    for i, (k, v) in enumerate(infos):
        bx = Inches(0.55) + i * Inches(3.2)
        add_rect(slide, bx, Inches(6.1), Inches(3.0), Inches(1.0),
                 fill=rgb_hex(0x1E,0x29,0x3B), line=C_PRIMARY, line_w=Pt(1))
        add_text(slide, k, bx + Inches(0.15), Inches(6.2), Inches(2.7), Inches(0.3),
                 font_size=Pt(9), color=C_MID_GRAY)
        add_text(slide, v, bx + Inches(0.15), Inches(6.48), Inches(2.7), Inches(0.45),
                 font_size=Pt(12), bold=True, color=C_WHITE)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — DAFTAR ISI
# ═════════════════════════════════════════════════════════════════════════════
def slide_daftar_isi(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "Daftar Isi", "Struktur presentasi proyek")

    items = [
        ("01", "Latar Belakang & Tujuan"),
        ("02", "Ruang Lingkup & Fitur Sistem"),
        ("03", "Arsitektur Sistem"),
        ("04", "Entity Relationship Diagram (ERD)"),
        ("05", "Logical Relational Schema (LRS)"),
        ("06", "Class Diagram"),
        ("07", "Activity Diagram"),
        ("08", "Sequence Diagram"),
        ("09", "Desain Antarmuka (UI)"),
        ("10", "Teknologi & Implementasi"),
        ("11", "Pengujian Sistem"),
        ("12", "Kesimpulan & Saran"),
    ]

    col1 = items[:6]
    col2 = items[6:]

    def draw_item(num, label, left, top):
        add_rect(slide, left, top, Inches(0.55), Inches(0.48),
                 fill=C_PRIMARY, line=None)
        add_text(slide, num, left, top, Inches(0.55), Inches(0.48),
                 font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, label, left + Inches(0.65), top + Inches(0.05),
                 Inches(5.2), Inches(0.38), font_size=Pt(13), color=C_WHITE)

    for i, (n, l) in enumerate(col1):
        draw_item(n, l, Inches(0.55), Inches(1.6) + i * Inches(0.85))
    for i, (n, l) in enumerate(col2):
        draw_item(n, l, Inches(6.9), Inches(1.6) + i * Inches(0.85))


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — LATAR BELAKANG
# ═════════════════════════════════════════════════════════════════════════════
def slide_latar_belakang(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "01  Latar Belakang", "Permasalahan yang mendorong pengembangan sistem")

    problems = [
        ("Pencatatan Manual", "Tabungan qurban masih dicatat di buku/Excel — rawan kesalahan dan sulit diaudit."),
        ("Transparansi Kurang", "Jamaah tidak bisa memantau saldo & status tabungan secara mandiri."),
        ("Koordinasi Kelompok", "Pengelompokan 7 anggota per sapi dilakukan manual tanpa sistem otomatis."),
        ("Pembayaran Tunai", "Setoran tunai sulit dilacak; tidak ada bukti digital yang tersentralisasi."),
    ]

    for i, (title, desc) in enumerate(problems):
        col = i % 2
        row = i // 2
        bx = Inches(0.5) + col * Inches(6.3)
        by = Inches(1.65) + row * Inches(2.45)
        add_rect(slide, bx, by, Inches(6.0), Inches(2.2),
                 fill=C_CARD_BG, line=C_PRIMARY, line_w=Pt(1))
        add_rect(slide, bx, by, Inches(0.08), Inches(2.2), fill=C_PRIMARY)
        add_text(slide, title, bx + Inches(0.25), by + Inches(0.18),
                 Inches(5.5), Inches(0.45), font_size=Pt(14), bold=True, color=C_PRIMARY_LT)
        add_text(slide, desc, bx + Inches(0.25), by + Inches(0.6),
                 Inches(5.5), Inches(1.4), font_size=Pt(12), color=C_MID_GRAY)

    # tujuan
    add_rect(slide, Inches(0.5), Inches(6.8), Inches(12.3), Pt(3), fill=C_ACCENT)
    add_text(slide,
             "Tujuan: Membangun sistem informasi berbasis web yang memudahkan pengelolaan tabungan qurban secara digital, transparan, dan terintegrasi dengan pembayaran online.",
             Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.55),
             font_size=Pt(11), color=C_MID_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — RUANG LINGKUP & FITUR
# ═════════════════════════════════════════════════════════════════════════════
def slide_fitur(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "02  Ruang Lingkup & Fitur Sistem",
                 "Fungsionalitas utama berdasarkan peran pengguna")

    roles = [
        ("Jamaah", C_PRIMARY, [
            "Registrasi & login akun",
            "Ajukan bergabung kelompok qurban",
            "Pantau saldo & progres tabungan",
            "Bayar via VA / QRIS / e-wallet",
            "Lihat kwitansi digital",
            "Notifikasi in-app & WhatsApp",
        ]),
        ("Admin Biasa", C_ACCENT, [
            "Kelola data jamaah",
            "Terima / tolak permintaan gabung",
            "Tempatkan jamaah ke kelompok",
            "Rekap transaksi harian",
            "Kirim notifikasi manual",
        ]),
        ("Kepala Admin", C_YELLOW, [
            "Semua akses Admin Biasa",
            "Buat & hapus kelompok qurban",
            "Kelola akun admin",
            "Laporan keuangan lengkap",
            "Pengaturan sistem",
        ]),
    ]

    col_w = Inches(4.1)
    for i, (role, color, feats) in enumerate(roles):
        bx = Inches(0.45) + i * Inches(4.25)
        by = Inches(1.6)
        add_rect(slide, bx, by, col_w, Inches(0.52), fill=color)
        add_text(slide, role, bx, by, col_w, Inches(0.52),
                 font_size=Pt(15), bold=True, color=C_DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, by + Inches(0.52), col_w, Inches(5.0),
                 fill=C_CARD_BG, line=color, line_w=Pt(1))
        for j, feat in enumerate(feats):
            fy = by + Inches(0.7) + j * Inches(0.72)
            add_rect(slide, bx + Inches(0.2), fy + Inches(0.15),
                     Inches(0.12), Inches(0.12), fill=color)
            add_text(slide, feat, bx + Inches(0.45), fy,
                     Inches(3.5), Inches(0.6), font_size=Pt(12), color=C_WHITE)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — ARSITEKTUR SISTEM
# ═════════════════════════════════════════════════════════════════════════════
def slide_arsitektur(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "03  Arsitektur Sistem",
                 "3-tier architecture: Presentation · Business Logic · Data")

    # tiga kotak tier
    tiers = [
        ("Presentation Tier", "React 18 + Vite\nTailwind CSS · Framer Motion\nReact Router DOM\nAxios · React Hook Form",
         C_PRIMARY),
        ("Business Logic Tier", "Node.js + Express.js\nJWT Authentication\nRBAC Middleware\nPayment Service\nWhatsApp Service",
         C_ACCENT),
        ("Data Tier", "MongoDB + Mongoose\n$jsonSchema Validators\nCollections:\njamaah · admin · sapi_qurban\ntransaksi · notifikasi",
         C_YELLOW),
    ]

    for i, (title, content, color) in enumerate(tiers):
        bx = Inches(0.45) + i * Inches(4.2)
        by = Inches(1.65)
        bw = Inches(3.9)
        bh = Inches(4.8)
        add_rect(slide, bx, by, bw, Inches(0.55), fill=color)
        add_text(slide, title, bx, by, bw, Inches(0.55),
                 font_size=Pt(13), bold=True, color=C_DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, by + Inches(0.55), bw, bh - Inches(0.55),
                 fill=C_CARD_BG, line=color, line_w=Pt(1))
        add_text(slide, content, bx + Inches(0.2), by + Inches(0.75),
                 bw - Inches(0.3), bh - Inches(0.9),
                 font_size=Pt(12), color=C_WHITE)

        # panah antar tier
        if i < 2:
            arr_x = bx + bw + Inches(0.08)
            arr_y = by + Inches(2.2)
            add_text(slide, "⟺", arr_x, arr_y, Inches(0.28), Inches(0.45),
                     font_size=Pt(18), color=C_MID_GRAY, align=PP_ALIGN.CENTER)

    # baris bawah: payment gateway
    add_rect(slide, Inches(0.45), Inches(6.55), Inches(12.35), Inches(0.72),
             fill=rgb_hex(0x1E,0x29,0x3B), line=C_PRIMARY, line_w=Pt(1))
    add_text(slide, "Payment Gateway Sandbox  ·  Webhook HMAC-SHA256  ·  Virtual Account / QRIS / E-Wallet  ·  Notifikasi WhatsApp (Fonnte API)",
             Inches(0.6), Inches(6.68), Inches(12.0), Inches(0.45),
             font_size=Pt(11), color=C_ACCENT, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — ERD PLACEHOLDER
# ═════════════════════════════════════════════════════════════════════════════
def slide_erd(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "04  Entity Relationship Diagram (ERD)",
                 "Relasi antar entitas pada basis data MongoDB")

    add_img_placeholder(slide,
                        Inches(0.5), Inches(1.6),
                        Inches(12.3), Inches(5.6),
                        label="[ Sisipkan ERD di sini ]",
                        sub_label="Entitas: Admin · Jamaah · Sapi Qurban · Kelompok Qurban · Tabungan · Detail Kelompok · Permintaan Gabung · Transaksi · Notifikasi",
                        bg=rgb_hex(0x1A,0x23,0x35),
                        border=C_PRIMARY,
                        label_color=C_PRIMARY_LT)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — LRS PLACEHOLDER
# ═════════════════════════════════════════════════════════════════════════════
def slide_lrs(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "05  Logical Relational Schema (LRS)",
                 "Pemetaan koleksi MongoDB ke skema relasional logis")

    add_img_placeholder(slide,
                        Inches(0.5), Inches(1.6),
                        Inches(12.3), Inches(5.6),
                        label="[ Sisipkan LRS di sini ]",
                        sub_label="Tabel: admin · jamaah · sapi_qurban · kelompok_qurban · tabungan_qurban · detail_kelompok · permintaan_gabung · transaksi · notifikasi",
                        bg=rgb_hex(0x1A,0x23,0x35),
                        border=C_ACCENT,
                        label_color=C_ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — CLASS DIAGRAM PLACEHOLDER
# ═════════════════════════════════════════════════════════════════════════════
def slide_class_diagram(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "06  Class Diagram",
                 "Struktur kelas model, controller, dan service")

    add_img_placeholder(slide,
                        Inches(0.5), Inches(1.6),
                        Inches(12.3), Inches(5.6),
                        label="[ Sisipkan Class Diagram di sini ]",
                        sub_label="Model · Controller · Service · Middleware · Utility",
                        bg=rgb_hex(0x1A,0x23,0x35),
                        border=C_YELLOW,
                        label_color=C_YELLOW)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — ACTIVITY DIAGRAM PLACEHOLDER (bagi per fitur)
# ═════════════════════════════════════════════════════════════════════════════
def slide_activity(prs):
    # Dua placeholder berdampingan: Registrasi & Tabungan/Pembayaran
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "07  Activity Diagram",
                 "Alur aktivitas proses utama sistem")

    activities = [
        ("Registrasi & Login", "Registrasi Jamaah · Verifikasi Akun · Login RBAC"),
        ("Setoran & Pembayaran", "Pilih Metode · Checkout · Webhook · Update Saldo"),
        ("Kelola Kelompok", "Buat Kelompok · Terima Anggota · Billing Aktif"),
    ]

    w = Inches(3.9)
    for i, (title, sub) in enumerate(activities):
        bx = Inches(0.45) + i * Inches(4.25)
        add_text(slide, title, bx, Inches(1.55), w, Inches(0.45),
                 font_size=Pt(12), bold=True, color=C_PRIMARY_LT, align=PP_ALIGN.CENTER)
        add_img_placeholder(slide,
                            bx, Inches(2.0), w, Inches(5.1),
                            label=f"[ Activity Diagram ]",
                            sub_label=sub,
                            bg=rgb_hex(0x1A,0x23,0x35),
                            border=C_PRIMARY,
                            label_color=C_MID_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — SEQUENCE DIAGRAM PLACEHOLDER
# ═════════════════════════════════════════════════════════════════════════════
def slide_sequence(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "08  Sequence Diagram",
                 "Interaksi antar objek untuk skenario utama")

    seqs = [
        ("Login & Autentikasi",      "Jamaah → Frontend → API → MongoDB → JWT"),
        ("Proses Pembayaran",        "Jamaah → Checkout → Payment GW → Webhook → DB"),
        ("Notifikasi Penempatan",    "Admin → Terima Permintaan → WhatsApp + In-App"),
    ]

    w = Inches(3.9)
    for i, (title, sub) in enumerate(seqs):
        bx = Inches(0.45) + i * Inches(4.25)
        add_text(slide, title, bx, Inches(1.55), w, Inches(0.45),
                 font_size=Pt(12), bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_img_placeholder(slide,
                            bx, Inches(2.0), w, Inches(5.1),
                            label=f"[ Sequence Diagram ]",
                            sub_label=sub,
                            bg=rgb_hex(0x1A,0x23,0x35),
                            border=C_ACCENT,
                            label_color=C_MID_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — USE CASE DIAGRAM PLACEHOLDER
# ═════════════════════════════════════════════════════════════════════════════
def slide_usecase(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "Diagram Tambahan — Use Case & Lainnya",
                 "Ruang untuk diagram tambahan sesuai kebutuhan")

    placeholders = [
        ("Use Case Diagram",   "Jamaah · Admin Biasa · Kepala Admin"),
        ("Diagram Lainnya",    "[ Tambahkan diagram sesuai kebutuhan ]"),
    ]

    for i, (title, sub) in enumerate(placeholders):
        bx = Inches(0.45) + i * Inches(6.5)
        add_text(slide, title, bx, Inches(1.55), Inches(6.0), Inches(0.45),
                 font_size=Pt(13), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_img_placeholder(slide,
                            bx, Inches(2.0), Inches(6.0), Inches(5.1),
                            label=f"[ {title} ]",
                            sub_label=sub,
                            bg=rgb_hex(0x1A,0x23,0x35),
                            border=C_MID_GRAY,
                            label_color=C_MID_GRAY)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — DESAIN UI (screenshot placeholders)
# ═════════════════════════════════════════════════════════════════════════════
def slide_ui(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "09  Desain Antarmuka (UI)",
                 "Tampilan halaman utama aplikasi web")

    screens = [
        ("Landing Page",       "Beranda publik dengan marquee pembayaran live"),
        ("Login / Register",   "Form 2-langkah wizard dengan password strength meter"),
        ("Dashboard Jamaah",   "Progres tabungan, riwayat transaksi, notifikasi"),
        ("Dashboard Admin",    "Rekap kelompok, manajemen jamaah, laporan"),
        ("Halaman Pembayaran", "Checkout VA / QRIS / e-wallet + kwitansi digital"),
        ("Manajemen Kelompok", "Status kelompok, daftar anggota, billing gate"),
    ]

    cols = 3
    for i, (title, desc) in enumerate(screens):
        col = i % cols
        row = i // cols
        bx = Inches(0.45) + col * Inches(4.25)
        by = Inches(1.6) + row * Inches(2.85)
        bw = Inches(4.0)
        bh = Inches(2.6)

        add_img_placeholder(slide, bx, by, bw, bh - Inches(0.55),
                            label=f"[ Screenshot: {title} ]",
                            sub_label="",
                            bg=rgb_hex(0x1A,0x23,0x35),
                            border=C_PRIMARY,
                            label_color=C_MID_GRAY)
        add_text(slide, title, bx, by + bh - Inches(0.6), bw, Inches(0.32),
                 font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, bx, by + bh - Inches(0.3), bw, Inches(0.28),
                 font_size=Pt(9), color=C_MID_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — TEKNOLOGI & IMPLEMENTASI
# ═════════════════════════════════════════════════════════════════════════════
def slide_teknologi(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "10  Teknologi & Implementasi",
                 "Stack lengkap yang digunakan dalam pengembangan")

    groups = [
        ("Frontend", C_PRIMARY, [
            "React 18.3 + Vite 8",
            "Tailwind CSS 3.4",
            "Framer Motion 12",
            "React Router DOM 6",
            "React Hook Form 7",
            "Axios 1.7",
            "React Toastify 10",
            "@heroicons/react 2",
        ]),
        ("Backend", C_ACCENT, [
            "Node.js + Express.js",
            "MongoDB + Mongoose 8",
            "JWT (jsonwebtoken)",
            "bcrypt 6",
            "Multer 2 + Cloudinary",
            "Express Validator",
            "Node-cron (scheduler)",
            "Dotenv",
        ]),
        ("Infrastruktur", C_YELLOW, [
            "MongoDB Atlas / Self-hosted",
            "Payment Gateway Sandbox",
            "Fonnte WhatsApp API",
            "HMAC-SHA256 Webhook",
            "Git + GitHub",
            "REST API Architecture",
            "RBAC Middleware",
            "JWT Bearer Token",
        ]),
    ]

    col_w = Inches(3.9)
    for i, (title, color, items) in enumerate(groups):
        bx = Inches(0.45) + i * Inches(4.25)
        by = Inches(1.65)
        add_rect(slide, bx, by, col_w, Inches(0.48), fill=color)
        add_text(slide, title, bx, by, col_w, Inches(0.48),
                 font_size=Pt(13), bold=True, color=C_DARK_BG, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, by + Inches(0.48), col_w, Inches(5.1),
                 fill=C_CARD_BG, line=color, line_w=Pt(1))
        for j, item in enumerate(items):
            iy = by + Inches(0.65) + j * Inches(0.58)
            add_rect(slide, bx + Inches(0.2), iy + Inches(0.17),
                     Inches(0.1), Inches(0.1), fill=color)
            add_text(slide, item, bx + Inches(0.42), iy,
                     col_w - Inches(0.55), Inches(0.48),
                     font_size=Pt(12), color=C_WHITE)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — PENGUJIAN SISTEM
# ═════════════════════════════════════════════════════════════════════════════
def slide_pengujian(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "11  Pengujian Sistem",
                 "Hasil pengujian fungsional dan keamanan")

    tests = [
        ("Autentikasi & Otorisasi", "✓", "JWT valid, expired token ditolak, RBAC sesuai peran"),
        ("Registrasi Pengguna",     "✓", "Validasi form, username unik, hash password bcrypt"),
        ("Pembayaran VA / QRIS",    "✓", "Checkout idempoten, webhook HMAC-SHA256 valid"),
        ("Anti-Double Transaction", "✓", "Duplikat transaksi pending terdeteksi & ditolak"),
        ("Billing Gate (7 anggota)","✓", "Tagihan hanya aktif ketika kelompok penuh (7/7)"),
        ("Notifikasi WhatsApp",     "✓", "Pesan terkirim via Fonnte API / mode simulasi"),
        ("Notifikasi In-App",       "✓", "Badge unread, polling 30 s, mark-read on open"),
        ("Marquee Real-time",       "✓", "Polling 5 s, antrian prioritas, format tanggal/jam"),
        ("Keamanan Input (XSS/SQLi)","✓","Express Validator + Mongoose schema sanitization"),
        ("Zero npm Vulnerabilities","✓", "npm audit clean setelah upgrade Vite 5 → 8"),
    ]

    headers = ["#", "Skenario Uji", "Status", "Keterangan"]
    col_widths = [Inches(0.4), Inches(3.8), Inches(0.9), Inches(7.0)]
    row_h = Inches(0.48)
    start_y = Inches(1.65)
    start_x = Inches(0.45)

    # header row
    x = start_x
    for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, start_y, cw, row_h, fill=C_PRIMARY)
        add_text(slide, hdr, x + Inches(0.05), start_y, cw, row_h,
                 font_size=Pt(11), bold=True, color=C_DARK_BG, align=PP_ALIGN.CENTER if j in (0,2) else PP_ALIGN.LEFT)
        x += cw

    for i, (scenario, status, notes) in enumerate(tests):
        y = start_y + (i + 1) * row_h
        row_fill = C_CARD_BG if i % 2 == 0 else rgb_hex(0x1A, 0x23, 0x35)
        x = start_x
        for j, (val, cw) in enumerate(zip([str(i+1), scenario, status, notes], col_widths)):
            add_rect(slide, x, y, cw, row_h, fill=row_fill,
                     line=rgb_hex(0x2D,0x3D,0x55), line_w=Pt(0.5))
            color = C_PRIMARY_LT if (j == 2 and val == "✓") else C_WHITE
            add_text(slide, val, x + Inches(0.05), y + Inches(0.04), cw, row_h - Inches(0.08),
                     font_size=Pt(10), color=color,
                     align=PP_ALIGN.CENTER if j in (0,2) else PP_ALIGN.LEFT)
            x += cw


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — KESIMPULAN & SARAN
# ═════════════════════════════════════════════════════════════════════════════
def slide_kesimpulan(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_header(slide, "12  Kesimpulan & Saran",
                 "Ringkasan hasil dan rekomendasi pengembangan lanjut")

    # kesimpulan
    add_rect(slide, Inches(0.45), Inches(1.6), Inches(5.8), Inches(0.5), fill=C_PRIMARY)
    add_text(slide, "Kesimpulan", Inches(0.45), Inches(1.6), Inches(5.8), Inches(0.5),
             font_size=Pt(14), bold=True, color=C_DARK_BG, align=PP_ALIGN.CENTER)
    points_k = [
        "Sistem berhasil mendigitalisasi pengelolaan tabungan qurban dari pencatatan manual ke platform web modern.",
        "Arsitektur 3-tier (React · Express · MongoDB) memastikan skalabilitas dan pemisahan tanggung jawab yang baik.",
        "Role-based access control (Jamaah · Admin Biasa · Kepala Admin) menjamin keamanan data antar pengguna.",
        "Integrasi payment gateway + notifikasi WhatsApp meningkatkan kepercayaan dan kemudahan jamaah.",
        "Antarmuka premium dengan animasi Framer Motion memberikan pengalaman pengguna kelas SaaS.",
    ]
    add_rect(slide, Inches(0.45), Inches(2.1), Inches(5.8), Inches(4.6),
             fill=C_CARD_BG, line=C_PRIMARY, line_w=Pt(1))
    for i, p in enumerate(points_k):
        add_rect(slide, Inches(0.6), Inches(2.28) + i * Inches(0.85),
                 Inches(0.1), Inches(0.1), fill=C_PRIMARY_LT)
        add_text(slide, p, Inches(0.8), Inches(2.2) + i * Inches(0.85),
                 Inches(5.3), Inches(0.78), font_size=Pt(11), color=C_WHITE)

    # saran
    add_rect(slide, Inches(6.55), Inches(1.6), Inches(6.3), Inches(0.5), fill=C_ACCENT)
    add_text(slide, "Saran Pengembangan", Inches(6.55), Inches(1.6), Inches(6.3), Inches(0.5),
             font_size=Pt(14), bold=True, color=C_DARK_BG, align=PP_ALIGN.CENTER)
    points_s = [
        "Integrasi payment gateway produksi (Midtrans / Xendit) untuk transaksi nyata.",
        "Fitur laporan & ekspor PDF/Excel untuk keperluan audit keuangan masjid.",
        "Progressive Web App (PWA) agar dapat diakses offline via smartphone.",
        "Notifikasi push browser menggunakan Web Push API.",
        "Fitur rekap otomatis saat Idul Adha mendekat dengan email reminder.",
        "Dashboard analitik dengan grafik tren tabungan per bulan/tahun.",
    ]
    add_rect(slide, Inches(6.55), Inches(2.1), Inches(6.3), Inches(4.6),
             fill=C_CARD_BG, line=C_ACCENT, line_w=Pt(1))
    for i, p in enumerate(points_s):
        add_rect(slide, Inches(6.7), Inches(2.28) + i * Inches(0.72),
                 Inches(0.1), Inches(0.1), fill=C_ACCENT)
        add_text(slide, p, Inches(6.9), Inches(2.2) + i * Inches(0.72),
                 Inches(5.8), Inches(0.65), font_size=Pt(11), color=C_WHITE)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — PENUTUP
# ═════════════════════════════════════════════════════════════════════════════
def slide_penutup(prs):
    slide = prs.slides.add_slide(BLANK)
    set_slide_bg(slide, C_DARK_BG)

    add_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=C_PRIMARY)
    add_rect(slide, Inches(0), SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill=C_PRIMARY)

    add_img_placeholder(slide,
                        Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5),
                        label="[ LOGO ]", bg=rgb_hex(0x1E,0x29,0x3B),
                        border=C_PRIMARY, label_color=C_PRIMARY_LT)

    add_text(slide, "Terima Kasih",
             Inches(0.55), Inches(2.2), Inches(9), Inches(1.3),
             font_size=Pt(56), bold=True, color=C_WHITE)
    add_text(slide, "Sistem Tabungan Qurban — Platform Digital Pengelolaan Qurban Masjid",
             Inches(0.55), Inches(3.5), Inches(9.5), Inches(0.6),
             font_size=Pt(16), color=C_PRIMARY_LT)

    add_rect(slide, Inches(0.55), Inches(4.1), Inches(4), Pt(3), fill=C_PRIMARY)

    add_text(slide, "Pertanyaan & Diskusi",
             Inches(0.55), Inches(4.3), Inches(6), Inches(0.5),
             font_size=Pt(14), bold=True, color=C_ACCENT)

    add_text(slide,
             "Repository: github.com/pentest-acc/Tabungan-Qurban\n"
             "Stack: React · Node.js · MongoDB · Tailwind CSS · Framer Motion",
             Inches(0.55), Inches(4.85), Inches(9), Inches(0.85),
             font_size=Pt(12), color=C_MID_GRAY)

    add_text(slide, f"© {datetime.date.today().year}  Sistem Tabungan Qurban",
             Inches(0.55), Inches(6.9), Inches(6), Inches(0.4),
             font_size=Pt(10), color=rgb_hex(0x47,0x55,0x69))


# ── BUILD ─────────────────────────────────────────────────────────────────────
slide_cover(prs)
slide_daftar_isi(prs)
slide_latar_belakang(prs)
slide_fitur(prs)
slide_arsitektur(prs)
slide_erd(prs)
slide_lrs(prs)
slide_class_diagram(prs)
slide_activity(prs)
slide_sequence(prs)
slide_usecase(prs)
slide_ui(prs)
slide_teknologi(prs)
slide_pengujian(prs)
slide_kesimpulan(prs)
slide_penutup(prs)

OUT = "Sistem_Tabungan_Qurban.pptx"
prs.save(OUT)
print(f"\n✅  File berhasil dibuat: {OUT}")
print(f"   Total slide: {len(prs.slides)}")
print("\nCara membuka:")
print("  - Windows/Mac : klik dua kali file .pptx")
print("  - Online      : upload ke Google Slides atau PowerPoint Online")
print("\nUntuk menyisipkan diagram:")
print("  Klik area abu-abu → Insert > Image (atau paste gambar diagram Anda)")
